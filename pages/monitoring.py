import streamlit as st
from PIL import Image
import page_service as ps
from pymongo import MongoClient, ASCENDING, DESCENDING
from datetime import datetime, timedelta
import pandas as pd
import plotly.express as px
from io import BytesIO

import streamlit_authenticator as stauth
import yaml
from yaml.loader import SafeLoader
import pydeck as pdk

from pptx import Presentation
from pptx.util import Inches
import tempfile

import requests
import base64
import io

from streamlit_sortables import sort_items
from streamlit_condition_tree import condition_tree

import random
import json


from pptx.dml.color import RGBColor

coordinates_data = {
    "KALIMANTAN BARAT": (-0.0263, 109.3425),
    "PAPUA SELATAN": (-5.4656, 140.6324),
    "DAERAH ISTIMEWA YOGYAKARTA": (-7.7972, 110.3688),
    "SULAWESI BARAT": (-2.6713, 119.4543),
    "BALI": (-8.4095, 115.1889),
    "ACEH": (4.6951, 96.7494),
    "SULAWESI TENGAH": (-1.4307, 121.4456),
    "PAPUA BARAT": (-1.3361, 133.1747),
    "KALIMANTAN TIMUR": (0.5387, 116.4194),
    "SULAWESI SELATAN": (-4.6796, 119.7320),
    "SULAWESI UTARA": (1.4379, 124.8489),
    "MALUKU UTARA": (1.6944, 127.8080),
    "KALIMANTAN UTARA": (2.1381, 117.4165),
    "JAWA TENGAH": (-7.1500, 110.1403),
    "PAPUA": (-4.2699, 138.0804),
    "SUMATERA SELATAN": (-3.3194, 104.9147),
    "GORONTALO": (0.6788, 122.4559),
    "PAPUA TENGAH": (-3.8825, 137.1627),
    "RIAU": (0.5104, 101.4381),
    "KALIMANTAN TENGAH": (-1.6815, 113.3824),
    "KEPULAUAN RIAU": (3.9457, 108.1429),
    "KALIMANTAN SELATAN": (-2.9741, 115.5557),
    "BANTEN": (-6.1202, 106.1503),
    "JAMBI": (-1.6108, 103.6119),
    "LAMPUNG": (-5.4500, 105.2663),
    "BENGKULU": (-3.8004, 102.2655),
    "SUMATERA BARAT": (-0.7371, 100.2711),
    "JAWA BARAT": (-6.9147, 107.6098),
    "NUSA TENGGARA BARAT": (-8.6529, 116.3249),
    "NUSA TENGGARA TIMUR": (-9.4750, 119.8707),
    "PAPUA PEGUNUNGAN": (-3.7330, 138.0883),
    "JAWA TIMUR": (-7.2504, 112.7688),
    "KEPULAUAN BANGKA BELITUNG": (-2.7455, 106.1134),
    "SULAWESI TENGGARA": (-4.1470, 122.1746),
    "SUMATERA UTARA": (2.2583, 99.6115),
    "MALUKU": (-3.2384, 129.4921),
    "DKI JAKARTA": (-6.2088, 106.8456)
}

# Inisiasi page
st.set_page_config(
    page_title="Analisa Data Pasker",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="auto"
)


color_sequence = px.colors.qualitative.D3

with open('config.yaml') as file:
    config_yaml = yaml.load(file, Loader=SafeLoader)

authenticator = stauth.Authenticate(
    config_yaml['credentials'],
    config_yaml['cookie']['name'],
    config_yaml['cookie']['key'],
    config_yaml['cookie']['expiry_days'],
    auto_hash=False
)

authenticator.login(location="unrendered")

if not st.session_state.authentication_status:
    st.switch_page("streamlit_app.py")

st.markdown("""
    <style> 
        [class*="st-key-card_dashoard_"]{
            padding: 15px 20px;
            border-radius: 5px;
            border: 2px solid #33CCCC;
            box-shadow: 3px 3px 1px 2px #33CCCC;
        }
        
        @media (max-width: 767px) {
            [class*="st-key-card_dashoard_"] > div > div:nth-child(2){
                display: none !important;
            }

            [class*="st-key-chart_card_"] div[data-baseweb='tab-list']{
                zoom: 0.5 !important;
            }
        }

        [class*="st-key-chart_card_"] div[data-baseweb='tab-list']{
            gap: 0px;
            justify-content: center;
            zoom: 0.8;
        }
            
        [class*="st-key-chart_card_"] button[data-baseweb='tab']{
            padding: 0px 50px;
            border: 2px solid gray;
            border-radius: 10px;
        }
            
        [class*="st-key-chart_card_"] button[data-baseweb='tab']:hover{
            border: 2px solid red;    
        }

        [class*="st-key-chart_card_"] button[aria-selected='true']{
            border: 2px solid red;
            background: #ffeded;
        } 
        
        [class*="st-key-chart_card_"] button[aria-selected='true'] *{
            font-weight: 900;
        }

        [class*="st-key-chart_card_"] button[data-baseweb='tab']:nth-child(1){
            border-top-right-radius: 0px;
            border-bottom-right-radius: 0px;
            # border-right: unset;
        }
            
        [class*="st-key-chart_card_"] button[data-baseweb='tab']:nth-child(2){
            border-radius: unset;
        }
            
        [class*="st-key-chart_card_"] button[data-baseweb='tab']:nth-child(3){
            border-top-left-radius: 0px;
            border-bottom-left-radius: 0px;
        }
            
        [class*="st-key-chart_card_"] div[data-baseweb='tab-highlight'], [class*="st-key-chart_card_"] div[data-baseweb='tab-border']{
            display: none;    
        }
            
        [class*="st-key-chart_card_"] .stNumberInput input{
            background: white;
            border-bottom: 1px solid gray;
        }
    </style>
""", unsafe_allow_html=True)
ps.setup_style_awal(st)

img_logo = Image.open('logo_pasker.png')
st.logo(img_logo, size="large")


if(authenticator):
    ps.setup_st_sidebar(st, authenticator)

with open("konfig.json", 'r') as file:
    config = json.load(file)

#init DB
@st.cache_resource
def init_connection():
    return MongoClient(**st.secrets["mongo"])

client = init_connection()
db = client.medmon
collection = db["hasil_proses_v1"]

# Main Page Function
# st.write("# Analisa Data 🕵️‍♂️🚀")
# st.subheader("Semua Data Hasil Scraping & Klasifikasi")

def get_coordinates(province_name):
    if province_name in coordinates_data:
        return coordinates_data[province_name]
    else:
        return None, None
    
original_kd = ["Jabatan", "Jabatan Detail", "Tipe Pekerjaan", "Tingkat Pekerjaan", "Tingkat Pendidikan", "Pengalaman Kerja", 
      "Tunjangan", "Jenis Kelamin", "Cara Kerja", "Lokasi", "Lokasi Kota", 
      "Keterampilan Bahasa", "Keterampilan Teknis", "Keterampilan Non Teknis", "Rentang Gaji", "Ukuran Perusahaan"]

if "kd" not in st.session_state:
    st.session_state["kd"] = original_kd
    
kd = st.session_state["kd"]

is_media_online = False

# Function to save figure and add to PowerPoint
def save_chart_to_slide(presentation, fig, title, df=None):
    slide_layout = presentation.slide_layouts[1]  # Title slide layout
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Save Plotly figure as an image
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
        fig.write_image(tmpfile.name, format="png", scale=2, width=1200, height=800)
        slide.shapes.add_picture(tmpfile.name, Inches(0.7), Inches(1.2), width=Inches(5.5))
    
    if df is not None and not df.empty:
        add_table_from_df_to_slide(presentation=presentation, df=df, title=None, slide=slide)

def add_table_from_df_to_slide(presentation, df, title=None, slide=None):
    if not slide:
        # Add a slide with the specified layout
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        x, y, cx, cy = Inches(1.2), Inches(1), Inches(7.5), Inches(4)
    else:
        x, y, cx, cy = Inches(6.5), Inches(1), Inches(2.7), Inches(3.1)

    if title:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

    # Define the number of rows and columns for the table
    rows, cols = df.shape[0] + 1, df.shape[1]
    
    # Add table placeholder (position and size: x, y, width, height)
    table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
    
    # Add header row
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = col_name
    
    # Add data to table
    for row_idx, (index, row) in enumerate(df.iterrows(), start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)

def ganti_text_di_ppt(slide, tag, data):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if '{{ ' + tag + ' }}' in text:
                shape.text_frame.text = text.replace('{{ ' + tag + ' }}', data)  
    

# Create PowerPoint Presentation
presentation = Presentation("template.pptx")
opening_slide = presentation.slides[0]


# Helper function to set the slider range based on the selectbox value
def get_slider_range(option):
    today = datetime.now()
    if option == "Hari ini":
        return [today.replace(hour=0, minute=0, second=0, microsecond=0), 
                today.replace(hour=23, minute=59, second=59, microsecond=0)]
    elif option == "Kemarin":
        yesterday = today - timedelta(days=1)
        return [yesterday.replace(hour=0, minute=0, second=0, microsecond=0), 
                yesterday.replace(hour=23, minute=59, second=59, microsecond=0)]
    elif option == "7 Hari Terakhir":
        last_7_days_start = today - timedelta(days=7)
        return [last_7_days_start.replace(hour=0, minute=0, second=0, microsecond=0), 
                today.replace(hour=23, minute=59, second=59, microsecond=0)]
    elif option == "1 Bulan Terakhir":
        last_month_start = today - timedelta(days=30)
        return [last_month_start.replace(hour=0, minute=0, second=0, microsecond=0), 
                today.replace(hour=23, minute=59, second=59, microsecond=0)]
    elif option == "3 Bulan Terakhir":
        last_3_months_start = today - timedelta(days=90)
        return [last_3_months_start.replace(hour=0, minute=0, second=0, microsecond=0), 
                today.replace(hour=23, minute=59, second=59, microsecond=0)]
    elif option == "Rentang Waktu":
        # Default range for custom range
        return [datetime(2024, 7, 1), datetime(2024, 11, 29)]
    return None

@st.cache_data(show_spinner=False)
def get_all_column():
    # result = collection.find_one().keys()
    return None

@st.cache_data(ttl=300, show_spinner=False)
def get_top_accounts(alur_waktu, is_media_online):
    match_stage = {
        "Klasifikasi Akun": "Akun Loker",
        "Tanggal Publikasi": {"$gte": alur_waktu[0], "$lte": alur_waktu[1]}
    }

    if not is_media_online:
        match_stage["Sumber"] = {"$ne": "Media Online"}

    pipeline = [
        {"$match": match_stage},  # Filter by category
        {
            "$group": {
                "_id": "$Akun/Judul",
                "Followers": {"$last": "$Followers"},
                "Sumber": {"$last": "$Sumber"}  # Include Sumber in the output
            }
        },
        {"$sort": {"Followers": -1}},
        {"$limit": 20},
        {
            "$project": {
                "_id": 1,
                "Sumber": 1,
                "Followers": 1,
            }
        }
    ]
    return pd.DataFrame(list(collection.aggregate(pipeline)))

@st.cache_data(ttl=300, show_spinner=False)
def get_total_posts_by_classification(alur_waktu, is_media_online):
    match_stage = { "Tanggal Publikasi": {"$gte": alur_waktu[0], "$lte": alur_waktu[1]}}

    if not is_media_online:
        match_stage["Sumber"] = {"$ne": "Media Online"}

    pipeline = [
        {"$match": match_stage},
        {"$group": {"_id": "$Klasifikasi Akun", "Count": {"$sum": 1}}}
    ]
    return pd.DataFrame(list(collection.aggregate(pipeline)))

@st.cache_data(ttl=300, show_spinner=False)
def get_category_counts(category_field,alur_waktu,is_media_online,custom_query=None):
    try:
        # If the category is "Jenis Kelamin", handle it specifically
        match_stage = {"$match": {}}
        match_stage["$match"]["Tanggal Publikasi"] = {"$gte": alur_waktu[0], "$lte": alur_waktu[1]}
        
        if not is_media_online:
            match_stage["$match"]["Sumber"] = {"$ne": "Media Online"}
        
        if custom_query and custom_query[1]:
            match_stage["$match"][custom_query[0]] = custom_query[1]

        if "conditional_tree_query" in st.session_state:
            conditional_tree_query = st.session_state["conditional_tree_query"]

            for key, value in conditional_tree_query.items():
                if value == "all":
                    # If the value is "all", do not add it to the match stage
                    continue
                else:
                    # If it's not "all", add the condition to $match
                    match_stage["$match"][key] = value

        if category_field == "Jenis Kelamin":
            pipeline = [
                match_stage,
                {
                    "$project": {
                        "Jenis Kelamin": {
                            "$cond": [
                                {
                                    "$eq": [
                                        {"$size": {"$setIntersection": [{"$ifNull": ["$Jenis Kelamin", []]}, ["Laki-Laki", "Perempuan"]]}},
                                        2
                                    ]
                                },
                                "Laki-laki & Perempuan",  # Both genders are present
                                {"$arrayElemAt": ["$Jenis Kelamin", 0]}  # Just one gender
                            ]
                        }
                    }
                },
                {"$group": {"_id": "$Jenis Kelamin", "Count": {"$sum": 1}}},
                {"$match": {"_id": {"$ne": "tdk_ada_informasi"}}}
            ]
        elif category_field == "Rentang Gaji":
            match_stage["$match"]["Digit Gaji (Clean)"] = {"$ne": None}
            
            pipeline = [
                match_stage,
                {
                    "$bucket": {
                        "groupBy": "$Digit Gaji (Clean)",  # Field to group by
                        "boundaries": [0, 3000000, 6000000, 10000000],  # Ranges
                        "default": "lebih dari 10jt",  # Label for outliers
                        "output": {
                            "Count": {"$sum": 1}  # Count documents in each bucket
                        }
                    }
                },
                {
                    "$project": {
                        "Rentang Gaji": {
                            "$switch": {  # Custom range labels
                                "branches": [
                                    {"case": {"$eq": ["$_id", 0]}, "then": "kurang dari 3jt"},
                                    {"case": {"$eq": ["$_id", 3000000]}, "then": "3jt - 5,9jt"},
                                    {"case": {"$eq": ["$_id", 6000000]}, "then": "6jt - 9,9jt"}
                                ],
                                "default": "lebih dari 10jt"
                            }
                        },
                        "Count": 1,
                        "_id": 0
                    }
                }
                
            ]
        else:
            # Default pipeline for other categories
            pipeline = [
                match_stage,
                {"$unwind": f"${category_field}"},
                {"$group": {"_id": f"${category_field}", "Count": {"$sum": 1}}},
                {"$match": {"_id": {"$ne": "tdk_ada_informasi"}}}
            ]
        # st.write(match_stage)
        
        # Execute the aggregation pipeline and return as DataFrame
        result = pd.DataFrame(list(collection.aggregate(pipeline)))

        
        if result.empty:
            result = pd.DataFrame(columns=[category_field, "Count"])

    except Exception as e:
        result = pd.DataFrame(columns=[category_field, "Count"])

    
    return result

@st.dialog("Download disini :")
def save_ppt():
    try:
        with st.spinner("Menyiapkan File PPT ke Google Drive ..."):
            data_presentation = st.session_state["presentation"]
            
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmpfile:
                data_presentation.save(tmpfile.name)

            # Read the file and encode in Base64
                with open(tmpfile.name, "rb") as file:
                    encoded_file = base64.b64encode(file.read()).decode("utf-8")

                # Prepare data payload
                data = {
                    "folderId": "1bMhZ6lnZQLO2hR2w73GQ1cEvD5VNsOnu",
                    "fileData": encoded_file,
                    "fileName": "dashboard_weekly_presentation_pasker.pptx"
                }

                # Endpoint URL
                endpoint = "https://script.google.com/macros/s/AKfycbypnLPxFKWlAFWcnx6U7OR28rsADxWBe5r2Np2UTjKMmQpQLYmzY0YiZw2lUDtTGh7_/exec?rute=upload_ppt"
                
                # Send POST request
                response = requests.post(endpoint, json=data)
        
        data_resp = response.json()
        url = data_resp.get('url', 'No URL returned')

        # st.write(data_resp)
        st.success("File Siap di Download", icon=":material/check_circle:")
        st.link_button("Lihat PPT", url, icon=":material/description:", type="primary", use_container_width=True)
    
    except:
        st.info("Mohon menunggu sampai data dashboard selesai ter-muat.", icon=":material/info:")

@st.dialog("Download disini :")
def save_excel():
    try:
        # fo = st.session_state["filter_option"]
        mv = st.session_state["max_value"]
        # st.write(fo)
        rd = st.slider("Rentang Data : ", value=[1,5000], min_value=1, max_value=mv)
        
        with st.spinner("Menyiapkan File Excel ..."):
            # df = fetch_data(search_query=None, sort_by="Tanggal Publikasi", sort_dir="🔼", filters=fo, page=rd[0], page_size=rd[1])
            df = None
            return
        
            buffer = io.BytesIO()
            with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
                df.to_excel(writer, index=False, sheet_name='Sheet1')

            buffer.seek(0)

            # Download button for the Excel file
            st.download_button(
                label="Download Excel File",
                data=buffer,
                file_name="mongodb_data.xlsx",
                type="primary",
                icon=":material/download:",
                use_container_width=True,
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
        # st.write(df)
    except:
        st.info("Mohon menunggu sampai data dashboard dan data scraping selesai ter-muat.", icon=":material/info:")

def rename_keys(data, key_mapping):
    return {key_mapping.get(key, key): value for key, value in data.items()}

def generate_default_condition_tree(fields):
    return {
        "type": "group",
        "properties": {
            "conjunction": "AND",
            "not": False
        },
        "children": [
            {
                "type": "rule",
                "properties": {
                    "fieldSrc": 'field',  # The 'fieldSrc' remains 'field' for dynamically generated rules
                    "field": field_name,
                    "operator": 'select_equals',  # Use '==' as the operator
                    "value": ["all"],  # Default value (first option)
                    "valueSrc": ['value'],  # Set to "value" to use the value directly
                    "valueType": ['select']  # Assuming "text" for simplicity
                }
            }
            for field_name, field_config in fields.items()
            if "fieldSettings" in field_config and "listValues" in field_config["fieldSettings"]
        ]
    }

@st.dialog("Sort & Filter Chart Data :", width="large")
def show_konfig_dashboard():
    # st.write("##### Urutan Data : ")
    # st.rerun()

    # st.json(st.session_state["data_config_filter"], expanded=False)

    if "tree" not in st.session_state:
        value = generate_default_condition_tree(config["fields"])
    else:
        value = st.session_state["tree"]


    with st.container(border=True):
        # col = st.columns([3,1.2])
        # with col[0]:
        # with  col[1]:

        # st.write("#### Tampilkan Data Media Online :")
        


        st.write("#### Urutan Data Chart :")
            
        revised_kd = sort_items(kd, direction="horizontal")

        st.write("#### Filter Data Chart :")
        query_tmp = condition_tree(
            config,
            tree=value,
            always_show_buttons=True,
            return_type='mongodb',
            placeholder="Tambahkan Rule Baru",
            key="tree"
        )

    cols = st.columns(3)

    with cols[2]:
        if st.button("Terapkan", use_container_width=True, type="primary", icon=":material/play_arrow:"):
            # Re-rename Data

            key_mapping = {
                "Jabatan": "Jabatan Detail",
                "Bidang Jabatan": "Jabatan",
                "Indikasi Scam": "Scam Detector",
            }

            query_tmp = rename_keys(query_tmp, key_mapping)

            st.session_state["kd"] = revised_kd
            st.session_state["conditional_tree_query"] = query_tmp
            
            st.cache_data.clear()
            st.rerun(scope='app')
        
def masukan_ke_data_config_filter(label, df):
    if "data_config_filter" not in st.session_state:
        st.session_state["data_config_filter"] = {
            'fields': {}
        }
    st.session_state["data_config_filter"]["fields"][label] = {
        "label": label,
        "type": "select",
        "fieldSettings": {
            "listValues": [
                {"value": row[label], "title": row[label]} for _, row in df.iterrows()
            ] + [{"value": "all", "title": "Semua Klasifikasi"}]  # Add additional option
        }
    }

def format_datetime_to_string(dt):
    return dt.strftime("%Y-%m-%d %H:%M:%S")

def get_chart_type(data_type,list_of_chart):
    if data_type == "Bidang Jabatan":
        return list_of_chart[2]  # Vertical Bar Chart
    elif data_type == "Jabatan":
        return list_of_chart[2]  # Vertical Bar Chart
    elif data_type == "Tipe Pekerjaan":
        return list_of_chart[0]  # Pie Chart
    elif data_type == "Tingkat Pekerjaan":
        return list_of_chart[0]  # Pie Chart
    elif data_type == "Jenjang Pendidikan":
        return list_of_chart[1]  # Horizontal Bar Chart
    elif data_type == "Pengalaman Kerja":
        return list_of_chart[1]  # Horizontal Bar Chart
    elif data_type == "Tunjangan":
        return list_of_chart[10]  # Donut Chart
    elif data_type == "Jenis Kelamin":
        return list_of_chart[10]  # Donut Chart
    elif data_type == "Model Kerja":
        return list_of_chart[10]  # Donut Chart
    elif data_type == "Provinsi":
        return list_of_chart[2]  # Vertical Bar Chart
    elif data_type == "Kota / Kabupaten":
        return list_of_chart[2]  # Vertical Bar Chart
    elif data_type == "Keterampilan Bahasa":
        return list_of_chart[1]  # Horizontal Bar Chart
    elif data_type == "Keterampilan Teknis":
        return list_of_chart[1]  # Horizontal Bar Chart
    elif data_type == "Keterampilan Non Teknis":
        return list_of_chart[1]  # Horizontal Bar Chart
    elif data_type == "Rentang Gaji":
        return list_of_chart[2]  # Vertical Bar Chart
    elif data_type == "Ukuran Perusahaan":
        return list_of_chart[2]  # Vertical Bar Chart
    elif data_type == "Indikasi Scam":
        return list_of_chart[5]  # Area Chart
    elif data_type == "Persepsi":
        return list_of_chart[2]  # Vertical Bar Chart
    else:
        return list_of_chart[0]

@st.fragment
def draw_chart(idx, item, aw, is_media_online):
    with st.container(border=True, key=f"chart_card_{idx}"):
        with st.spinner('Preparing data...'):
            cols = st.columns([6,1])

            with cols[1]:
                popover = st.popover("", icon=":material/filter_alt:")

            data_selected = popover.selectbox("Pilih Data :", config["fields"], key=f"popover_kd{idx}")
            arr_value = [item["title"] for item in config["fields"][data_selected]["fieldSettings"]["listValues"]]
            value_selected = popover.multiselect("Nilai Data :", arr_value, key=f"popover_kd_value{idx}")
            custom_query = [data_selected, value_selected]

            category_count = get_category_counts(item, aw, is_media_online,custom_query)

            # Rename Data
            match item:
                case "Jabatan":
                    item = "Bidang Jabatan"
                case "Jabatan Detail":
                    item = "Jabatan"
                case "Scam Detector":
                    item = "Indikasi Scam"
                case "Tingkat Pendidikan":
                    item = "Jenjang Pendidikan"
                case "Cara Kerja":
                    item = "Model Kerja"
                case "Lokasi":
                    item = "Provinsi"
                case "Lokasi Kota":
                    item = "Kota / Kabupaten"
                case _:
                    item = item

            with cols[0]:
                st.write(f'##### Visualisasi Data :red[{item}] di Setiap Postingan')
            
            if category_count.empty:
                   st.write(f"##### Data :red[{item}] Tidak Tersedia")
                   st.write(f"Data {item} tidak tesedia berdasarkan Filter yang dipilih.")
                   st.image("https://cdn.vectorstock.com/i/500p/04/13/no-data-empty-concept-vector-41830413.jpg", use_container_width=True)

                   return

            tab_chart = st.tabs(["Chart", "Data", "Settings"])
            
            with tab_chart[2]:
                # category_count.columns = [item, 'Count']
                category_count.columns = [item, 'Count'] if item != "Rentang Gaji" else ['Count', item]

                cols = st.columns(2)

                with cols[0]:
                    list_of_chart = [
                        "Pie Chart",
                        "Horizontal Bar Chart",
                        "Vertical Bar Chart",
                        "Scatter Plot",
                        "Line Chart",
                        "Area Chart",
                        "Funnel Chart",
                        "Treemap",
                        "Radar Chart",
                        "Sunburst Chart",
                        "Donut Chart"
                    ]

                    
                    chart_type = get_chart_type(item, list_of_chart)
                    # Add dropdown for chart type selection

                    default_idx = list_of_chart.index(chart_type)

                    chart_type = st.selectbox(
                        "Ganti Tipe Chart :",
                        list_of_chart,
                        index=default_idx,
                        placeholder="Ganti Tipe Chart disini ...",
                        key=f"chart_type_{idx}"
                    )
    
                with cols[1]:
                    sort_chart = st.selectbox(
                        "Urutkan Data Chart:",
                        ["Jumlah Paling Banyak", "Jumlah Paling Sedikit", "Urutkan Manual"] if len(category_count) <= 10 else ["Jumlah Paling Banyak", "Jumlah Paling Sedikit"],
                        key=f"chart_sort_{idx}"
                    )
                
                
                # Let the user define a custom order using streamlit_sortables
                if sort_chart == "Urutkan Manual":
                    st.write("Urutkan data:")
                    default_order = list(category_count[item])
                    custom_order = sort_items(default_order)
                    category_count = category_count.set_index(item).loc[custom_order].reset_index()
                else:
                    st.write("Tampilkan Data:")
                    default_count = category_count
                    col = st.columns([1,6], gap="small", vertical_alignment="bottom")

                    with col[0]:
                        data_awal = 1
                        
                        if len(category_count) <=9:
                            data_awal = len(category_count)
                        elif len(category_count) > 9 and len(category_count)  <= 60:
                            data_awal = round(len(category_count) / 3)
                        else:
                            data_awal = 25

                        data_limit = st.number_input(
                            "",
                            min_value=1,
                            max_value=len(category_count),
                            value=data_awal,
                            step=10,
                            key=f"data_limit_{idx}",
                            label_visibility="collapsed"
                        )

                        category_count = category_count.sort_values(by='Count', ascending=True if sort_chart=="Jumlah Paling Sedikit" else False)
                        category_count = category_count.head(data_limit)
                    with col[1]:
                        st.write(f"{item} {sort_chart}, dari :red[{len(default_count)}] Data.")
                    # st.radio("Tampilkan Data :", [f"10 {item} {sort_chart}",f"15 {item} {sort_chart}",f"20 {item} {sort_chart}",f"25 {item} {sort_chart}",f"30 {item} {sort_chart}"])
                    # category_count = category_count.head(data_limit)

            # chart_type = random.choice(list_of_chart)
            # chart_type = list_of_chart[idx % len(list_of_chart)] if idx > 0 and chart_type=="Pie Chart" else chart_type
            # Tabs for chart and data

            with tab_chart[0]:  # Chart tab
                if chart_type == "Pie Chart":
                    fig = px.pie(
                        category_count,
                        names=item,
                        values='Count',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Horizontal Bar Chart":
                    fig = px.bar(
                        category_count,
                        x='Count',
                        y=item,
                        orientation='h',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Vertical Bar Chart":
                    fig = px.bar(
                        category_count,
                        y='Count',
                        x=item,
                        orientation='v',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Scatter Plot":
                    fig = px.scatter(
                        category_count,
                        x=item,
                        y='Count',
                        color=item,
                        size='Count',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Line Chart":
                    fig = px.line(
                        category_count,
                        x=item,
                        y='Count',
                        markers=True,
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Area Chart":
                    fig = px.area(
                        category_count,
                        x=item,
                        y='Count',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Funnel Chart":
                    fig = px.funnel(
                        category_count,
                        y=item,
                        x='Count',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Treemap":
                    fig = px.treemap(
                        category_count,
                        path=[item],
                        values='Count',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Radar Chart":
                    df_radar = category_count.pivot_table(index=item, values='Count').reset_index()
                    fig = px.line_polar(
                        df_radar,
                        r='Count',
                        theta=item,
                        line_close=True
                    )
                elif chart_type == "Sunburst Chart":
                    fig = px.sunburst(
                        category_count,
                        path=[item],
                        values='Count',
                        color_discrete_sequence=color_sequence
                    )
                elif chart_type == "Donut Chart":
                    fig = px.pie(
                        category_count,
                        names=item,
                        hole=0.5,
                        values='Count',
                        color_discrete_sequence=color_sequence
                    )
                st.plotly_chart(fig)
                save_chart_to_slide(presentation, fig, f'Distribusi Data {item}', category_count)

            with tab_chart[1]:  # Data tab
                st.dataframe(category_count, hide_index=True, use_container_width=True)

    # masukan_ke_data_config_filter(item, default_count)
    # Add spacing after the chart
    st.text("")
    st.text("")

def convert_to_currency(value):
    """
    Convert a numeric value to Indonesian currency notation.

    Parameters:
        value (int or float): The numeric value to convert.

    Returns:
        str: The converted value with 'rb' or 'jt' as a suffix.
    """
    if value < 1000:
        return str(value)  # Return the value as is if it's less than 1000.

    elif value < 1_000_000:
        # Convert to 'rb'
        return f"{value / 1000:.0f} rb"

    else:
        # Convert to 'jt'
        return f"{value / 1_000_000:.1f} jt".replace('.', ',')

def generate_leaderboard_html(dataframe):
    leaderboard_html = """
        <style>
            .leaderboard-container {
                font-family: Arial, sans-serif;
                width: 100%;
                height: 500px;
                overflow-y: auto;
                padding: 10px 10px;
            }
            .leaderboard-item {
                display: flex;
                align-items: center;
                background: white;
                padding: 10px;
                margin-bottom: 10px;
                border-radius: 6px;
                box-shadow: 0 0px 4px rgba(0, 0, 0, 0.5);
                transition: transform 0.2s ease-in-out;
                text-decoration: none;
            }
            .leaderboard-item:hover {
                transform: translateY(-5px);
            }
            .rank {
                font-size: 20px;
                font-weight: bold;
                width: 40px;
                text-align: center;
                color: #4a5568;
            }
            .avatar {
                width: 50px;
                height: 50px;
                border-radius: 100%;
                margin-right: 15px;
                border: 2px solid #e2e8f0;
            }
            .info {
                flex: 1;
            }
            .name {
                font-size: 18px;
                font-weight: bold;
                margin: 0;
                color: black;

                overflow: hidden;
                display: -webkit-box;
                -webkit-line-clamp: 2; /* number of lines to show */
                        line-clamp: 2; 
                -webkit-box-orient: vertical;
            }
            .followers {
                color: #4a5568;
                margin: 5px 0 0;
                
                padding: 2px 8px;
                background: #FF7F0E;
                color: white;
                border-radius: 4px;
                font-size: 12px;
                display: inline-block;
            }
            .platform {
                display: inline-block;
                margin-top: 5px;
                padding: 2px 8px;
                background: #1F77B4;
                color: white;
                border-radius: 4px;
                font-size: 12px;
            }
        </style>
        <div class="leaderboard-container">
    """
    
    for i, row in dataframe.iterrows():
        leaderboard_html += f"""
        <div class="leaderboard-item">
            <div class="rank">#{i + 1}</div>
            <img src="https://cdn0.iconfinder.com/data/icons/social-messaging-ui-color-shapes/128/user-male-circle-blue-512.png" alt="Profile Picture" class="avatar">
            <div class="info">
                <div class="name">{row['_id']}</div>
                <span class="followers">{convert_to_currency(row['Followers'])} Followers</span>
                <span class="platform">{row['Sumber']}</span>
            </div>
        </div>
        """
    leaderboard_html += "</div>"

    return leaderboard_html

@st.fragment
def fragment_button_konfigurasi():
    if st.button("Sort & Filter Data", type="primary", icon=":material/tune:", use_container_width=True, disabled=True):
        show_konfig_dashboard()

with st.container(border=True):
    col = st.columns([3.5,1,1], vertical_alignment="center")

    with col[0]:
        st.write("## Dashboard Analisa Data 🕵️‍♂️🚀")
        st.write("###### Menampilkan Semua Data Hasil Klasifikasi Data Scraping Sosial Media & Media Online")
    with col[1]:
        popover = st.popover("Export Data", icon=":material/download:", use_container_width=True)

        if popover.button("Dashboard (PPT)", icon=":material/animated_images:", use_container_width=True):
            save_ppt()

        if popover.button("Scraping (XLSX)", icon=":material/table:", use_container_width=True, disabled=True):
            save_excel()
    with col[2]:
        fragment_button_konfigurasi()
        

    main_tabs = st.tabs(["Lowongan Pekerjaan", "Persepsi & Sentimen Netizen"])
    
    with main_tabs[0]:
        col = st.columns([3.3,0.7,1], vertical_alignment="center")

        default_start_date = datetime(2024, 7, 1)
        default_end_date = datetime(2024, 12, 31)

        with col[0]:
            st.write("#### Ringkasan Data")
        with col[1]:
            is_media_online = st.toggle("Media Online", value=is_media_online)
        with col[2]:
            # Use st.date_input for a date range
            col_child = st.columns([1, 10], vertical_alignment="center", gap="small")

            with col_child[0]:
                st.write("📆")
            with col_child[1]:
                alur_waktu = st.date_input(
                    "Select Date Range",
                    value=(default_start_date, default_end_date),
                    min_value=datetime(2024, 1, 1),
                    max_value=datetime(2024, 12, 31),
                    label_visibility="collapsed"
                )

            

        # Ensure valid date range
            if len(alur_waktu) == 2 and alur_waktu[0] > alur_waktu[1]:
                st.error("Start date cannot be after the end date.")

        aw = [
            date if isinstance(date, str) else date.strftime("%Y-%m-%d")
            for date in alur_waktu
        ]

        if len(aw) < 2:
            st.toast("Tanggal Harus memiliki 2 rentang waktu.")
            st.stop()

        tgl_ppt = " s/d ".join(str(x) for x in aw)
        ganti_text_di_ppt(opening_slide, "tanggal_data", f"({ tgl_ppt })")
        
        col = st.columns(3, vertical_alignment="center", gap="medium")

        with col[0]:
            with st.container(key="card_dashoard_1"):
                start_date, end_date = aw[0], aw[1]

                query = {
                    "Tanggal Publikasi": {
                        "$gte": start_date,
                        "$lt": end_date
                    }
                }

                if not is_media_online:
                    query["Sumber"] = {"$ne": "Media Online"}
                
                data_count_in_range = collection.count_documents(query)

                # Query total data up to the start and end dates for comparison
                total_data_up_to_start = collection.count_documents({
                    "Tanggal Publikasi": {"$lt": start_date}
                })

                total_data_up_to_end = collection.count_documents({
                    "Tanggal Publikasi": {"$lt": end_date}
                })

                # Calculate differences
                difference = total_data_up_to_end - total_data_up_to_start

                if total_data_up_to_start > 0:
                    percentage_difference = round((difference / total_data_up_to_start) * 100)
                else:
                    percentage_difference = 0  # Avoid division by zero

                col_child = st.columns([2.5,1])

                with col_child[0]:
                    st.metric("Total Data Scraping", f"{data_count_in_range:,}".replace(",", "."), delta=f"{percentage_difference}%")

                with col_child[1]:
                    st.image("https://cdn-icons-gif.flaticon.com/8112/8112604.gif", use_container_width=True)
        with col[1]:
            with st.container(key="card_dashoard_2"):
                query_conditions = [
                    {"Persentase Lowongan": {"$gt": 20}},
                    {"Tanggal Publikasi": {"$gte": aw[0], "$lte": aw[1]}}
                ]

                if not is_media_online:
                    query_conditions.append({"Sumber": {"$ne": "Media Online"}})

                total_lowongan = collection.count_documents({"$and": query_conditions})

                col_child = st.columns([2.5,1])

                with col_child[0]:
                    st.metric("Total Job Posting", f"{total_lowongan:,}".replace(",", "."), delta="0 %")
                with col_child[1]:
                    st.image("https://cdn-icons-gif.flaticon.com/6172/6172508.gif", use_container_width=True)
        with col[2]:
            with st.container(key="card_dashoard_3"):
                match_stage = {"$match": {"Persentase Lowongan": {"$gt": 20}, "Tanggal Publikasi": {"$gte": aw[0], "$lte": aw[1]}}}

                if not is_media_online:
                    match_stage["$match"]["Sumber"] = {"$ne": "Media Online"}

                pipeline = [
                    match_stage,
                    {"$group": {"_id": None, "totalQuota": {"$sum": "$Digit Kouta (Clean)"}}}
                ]

                # Menjalankan pipeline
                result = list(collection.aggregate(pipeline))

                if len(result) <= 0:
                    total_quota = 0
                else:
                    total_quota = result[0]["totalQuota"]

                col_child = st.columns([2.5,1])
                
                with col_child[0]:
                    st.metric("Total Kuota Pekerjaan", f"{round(total_quota):,}".replace(",", "."), delta="0%")
                with col_child[1]:
                    st.image("https://cdn-icons-gif.flaticon.com/7211/7211849.gif", use_container_width=True)

        with st.container(border=True):
            col = st.columns([5,1.2], vertical_alignment="center")
            
            with col[0]:
                st.write("##### Distribusi Jumlah Lowongan Berdasarkan :red[Lokasi Provinsi]")
            with col[1]:
                format_map = st.radio("", ["Map 2D", "Map 3D"], horizontal=True, label_visibility="collapsed")

            df_lokasi = get_category_counts("Lokasi",aw, is_media_online, None)
            df_lokasi.columns = ["Lokasi", 'Count']

            if df_lokasi.empty:
                df_lokasi = pd.DataFrame(columns=["Lokasi", "Count", "lat", "lon"])
            else:
                # Add coordinates to the DataFrame
                df_lokasi[['lat', 'lon']] = df_lokasi['Lokasi'].apply(lambda x: pd.Series(get_coordinates(x)))

                # Filter out any rows with missing coordinates
                df_lokasi = df_lokasi.dropna(subset=['lat', 'lon'])
            
            df_lokasi['Size'] = df_lokasi['Count'] * 5  # Scale the size by multiplying by 100
            df_lokasi['Elevation'] = df_lokasi['Count'] * 5  # Scale the elevation by multiplying by 100
            

            if format_map=="Map 2D":
                st.map(
                    df_lokasi,
                    latitude="lat",
                    longitude="lon",
                    size="Size",  # Use the 'Count' column to scale the dot size
                    zoom=4
                )
            else:
                layer = pdk.Layer(
                    "ColumnLayer",  # Use ColumnLayer for 3D columns
                    data=df_lokasi,
                    get_position='[lon, lat]',
                    get_elevation='Elevation',  # Use 'Elevation' column for 3D height
                    elevation_scale=10,  # Adjust this scale to modify column height
                    radius=50000,  # Adjust radius for column size
                    get_fill_color='[255, 0, 0, 150]',  # RGBA color for columns
                    pickable=True,
                    auto_highlight=True
                )

                # Define the view state for the map
                view_state = pdk.ViewState(
                    latitude=df_lokasi['lat'].mean(),
                    longitude=df_lokasi['lon'].mean(),
                    zoom=4,
                    pitch=50,  # Angle for 3D effect
                    bearing=0
                )

                # Add tooltip configuration
                tooltip = {
                    "html": "<b>Province:</b> {Lokasi} <br/> <b>Job Count:</b> {Count}",
                    "style": {"color": "white"}
                }

                # Create the PyDeck chart
                st.pydeck_chart(pdk.Deck(
                    map_style="mapbox://styles/mapbox/light-v9",
                    initial_view_state=view_state,
                    layers=[layer],
                    tooltip=tooltip
                ))
        
        top_chart = st.columns(3, gap="small")

        if data_count_in_range <= 0:
            st.toast("Data Tidak Tersedia, di rentang waktu tsb.")
            st.stop()
        
        # Top 10 Accounts with Most Followers
        with top_chart[0]:
            with st.container(border=True):
                with st.spinner('Preparing data.'):
                    st.write("##### Top 20 :red[Akun Loker Follower] Terbanyak")
                    follower_by_account = get_top_accounts(aw,is_media_online)
                    # st.dataframe(follower_by_account, hide_index=True, use_container_width=True)

                    st.markdown(generate_leaderboard_html(follower_by_account), unsafe_allow_html=True)

                    add_table_from_df_to_slide(presentation=presentation, slide=None, df=follower_by_account, title="Top 20 Akun Loker Follower Terbanyak")

        # Total Posts by Sumber Media Sosial
        with top_chart[1]:
            with st.container(border=True):
                with st.spinner('Preparing data.'):
                    st.write("##### Visualisasi Data Berdasarkan :red[Sumber Sosial Media]")
                    
                    match_stage = {
                        "Tanggal Publikasi": {"$gte": aw[0], "$lte": aw[1]}
                    }

                    if not is_media_online:
                        match_stage["Sumber"] = {"$ne": "Media Online"}

                    # Query to count each source
                    source_count = collection.aggregate([
                        {"$match": match_stage},
                        {"$group": {"_id": "$Sumber", "Count": {"$sum": 1}}}
                    ])
                    source_count_df = pd.DataFrame(list(source_count))
                    source_count_df.columns = ['Sumber', 'Count']

                    # Pie chart for source distribution
                    fig_source_pie = px.pie(source_count_df, names='Sumber', values='Count', color_discrete_sequence=color_sequence)
                    st.plotly_chart(fig_source_pie)
                    save_chart_to_slide(presentation, fig_source_pie, "Distribusi Sumber Sosial Media", source_count_df)

                    # masukan_ke_data_config_filter("Sumber", source_count_df)

        # Total Posts by Account Classification
        with top_chart[2]:
            with st.container(border=True):
                with st.spinner('Preparing data.'):
                    st.write("##### Total Postingan Berdasarkan :red[Klasifikasi Akun]")
                    sentimen_count = get_total_posts_by_classification(aw,is_media_online)
                    sentimen_count.columns = ['Klasifikasi Akun', 'Count']
                    fig_sentimen = px.bar(sentimen_count, x='Klasifikasi Akun', y='Count', color_discrete_sequence=color_sequence)
                    st.plotly_chart(fig_sentimen)
                    save_chart_to_slide(presentation, fig_sentimen, "Total Postingan Berdasarkan Klasifikasi Akun", sentimen_count)

                    # masukan_ke_data_config_filter("Klasifikasi Akun", sentimen_count)
        
        st.text("")
        st.text("")

        kolom_loop = st.columns(2, gap="medium")
        # Loop through categories and generate charts
        for idx, item in enumerate(kd):
            if idx % 2 == 0:
                with kolom_loop[0]:
                    draw_chart(idx,item,aw,is_media_online)
            else:
                with kolom_loop[1]:
                    draw_chart(idx,item,aw,is_media_online)
        
        
    with main_tabs[1]:
        kd_netizen = ["Sentimen", "Persepsi"]

        kolom_loop = st.columns(2, gap="medium")
        
        with kolom_loop[0]:
            draw_chart(len(kd), kd_netizen[0], aw, is_media_online)
        with kolom_loop[1]:
            draw_chart(len(kd) + 1, kd_netizen[1], aw, is_media_online)


        if "presentation" not in st.session_state:
            st.session_state["presentation"] = None

        st.session_state["presentation"] = presentation